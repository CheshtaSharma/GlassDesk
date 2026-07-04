import os
import json
import re
import csv
import threading
import uuid
import time
import random
from io import BytesIO, StringIO

import numpy as np
import faiss
import pyodbc
from flask import Flask, request, jsonify, render_template, Response, stream_with_context
from pypdf import PdfReader
from dotenv import load_dotenv
import google.generativeai as genai

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "dev")
app.config["MAX_CONTENT_LENGTH"] = int(os.getenv("MAX_UPLOAD_MB", "150")) * 1024 * 1024

genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
MODEL = os.getenv("GEMINI_MODEL", "gemini-1.5-flash")
EMBED_MODEL = "models/text-embedding-004"

CHUNK_SIZE = 900
MAX_TURNS = 20
EMBED_BATCH = 20

TOOLS = [
    {
        "function_declarations": [
            {
                "name": "get_index",
                "description": "Get a chapter-level overview of the document (chapter number, title, page range, chunk count). Call this first to understand document structure. For older documents with no chapter data, returns a flat index of every chunk instead.",
                "parameters": {"type": "object", "properties": {}},
            },
            {
                "name": "get_chapter_chunks",
                "description": "Get the chunk index (id, page, short preview) for every chunk inside one chapter. Use the chapter number returned by get_index. Only needed when get_index returned chapters.",
                "parameters": {
                    "type": "object",
                    "properties": {"chapter_number": {"type": "integer"}},
                    "required": ["chapter_number"],
                },
            },
            {
                "name": "get_chunk",
                "description": "Fetch the full text of one specific chunk by its id (e.g. 'c14').",
                "parameters": {
                    "type": "object",
                    "properties": {"chunk_id": {"type": "string"}},
                    "required": ["chunk_id"],
                },
            },
            {
                "name": "search_chunks",
                "description": "Keyword search across all chunks. Tries multiple normalizations automatically. Use broad keywords if specific ones fail.",
                "parameters": {
                    "type": "object",
                    "properties": {"keyword": {"type": "string"}},
                    "required": ["keyword"],
                },
            },
            {
                "name": "semantic_search",
                "description": "Meaning-based FAISS vector search. Use this when keyword search fails or for conceptual questions.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {"type": "string"},
                        "top_k": {"type": "integer"},
                    },
                    "required": ["query"],
                },
            },
        ]
    }
]

SYSTEM_PROMPT = (
    "You are answering a question about a document you do NOT have full access to. "
    "You have five tools: get_index, get_chapter_chunks, get_chunk, search_chunks, and semantic_search.\n\n"
    "This document may be split into chapters. get_index returns either a chapter-level "
    "overview (chapter number, title, page range, chunk count) or, for older documents "
    "with no chapter data, a flat index of every chunk. If get_index returns chapters, "
    "call get_chapter_chunks(chapter_number) for the chapter that looks relevant to the "
    "question before fetching individual chunks — do not assume chunk IDs from one "
    "chapter apply to another.\n\n"
    "STRICT RULES:\n"
    "1. NEVER guess or invent chunk IDs. Only use chunk IDs you have actually seen "
    "returned by get_index, get_chapter_chunks, search_chunks, or semantic_search in "
    "THIS conversation.\n"
    "2. If get_index shows chapters, prefer narrowing to the right chapter with "
    "get_chapter_chunks first. Otherwise (or if the topic could be anywhere), use "
    "search_chunks or semantic_search directly to find relevant chunks. "
    "Then call get_chunk on the IDs returned. Never call get_chunk with an ID you haven't "
    "seen in a tool result.\n"
    "3. If search_chunks returns no results, try a shorter or simpler keyword — for example "
    "if 'electrical resistivity' fails, try just 'resistivity' or 'resistance'.\n"
    "4. If one search returns nothing useful, try at least 2-3 different keyword variations "
    "before concluding the topic isn't in the document.\n"
    "5. If get_chunk returns an error, do NOT stop — go back and search again with a "
    "different keyword to find the correct chunk ID.\n"
    "6. Read multiple chunks if the topic spans several pages — fetch all relevant ones "
    "before answering. When a topic looks substantial, also fetch neighboring chunks "
    "(e.g. c14 and c15, or the previous/next page) so your explanation has full context "
    "rather than stopping mid-explanation.\n\n"
    "Once you have enough information, write your FINAL ANSWER (no tool call) as "
    "structured study notes, not a summary paragraph. The reader wants to actually learn "
    "the topic. Use markdown — it will be rendered nicely, so use it properly:\n"
    "- '## ' for section headings, '### ' for sub-headings, '**text**' for key terms and "
    "important values, '- ' for bullet lists, '1. ' for numbered steps.\n\n"
    "Follow this structure:\n"
    "## Quick Answer\n"
    "One to two sentences that directly answer the question.\n\n"
    "## Details\n"
    "Break the topic into logical sub-headings (###) if it has multiple parts. Under each: "
    "explain the concept fully in your own words — definitions (**bold** the term itself), "
    "the underlying reasoning or process (as numbered steps if it's a procedure or "
    "sequence), relevant facts/figures, and any worked examples or sample values from the "
    "document, described in full. Do not compress a multi-step explanation into one line — "
    "walk through it. Define technical terms the first time you use them.\n\n"
    "## Sources\n"
    "A short bullet list of the page(s) and chunk id(s) you drew from, e.g. "
    "'- p. 12, c14 — definition of resistance'.\n\n"
    "Other rules:\n"
    "- If different parts of the document disagree, note the discrepancy in Details.\n"
    "- Prefer thorough over brief — a longer, complete answer beats a short one, as long "
    "as everything is grounded in the document.\n"
    "- If you truly cannot find the answer after trying multiple searches, skip the above "
    "structure and say so honestly, listing exactly which keywords you tried."
)


# ---------------------------------------------------------------------------
# Database
# ---------------------------------------------------------------------------

def get_db():
    server = os.getenv("SQL_SERVER")
    database = os.getenv("SQL_DATABASE")
    driver = os.getenv("SQL_DRIVER", "ODBC Driver 17 for SQL Server")
    trusted = os.getenv("SQL_TRUSTED_CONNECTION", "false").lower() == "true"
    if trusted:
        conn_str = f"DRIVER={{{driver}}};SERVER={server};DATABASE={database};Trusted_Connection=yes;"
    else:
        user = os.getenv("SQL_USERNAME")
        pwd = os.getenv("SQL_PASSWORD")
        conn_str = f"DRIVER={{{driver}}};SERVER={server};DATABASE={database};UID={user};PWD={pwd};"
    return pyodbc.connect(conn_str, autocommit=False)


# ---------------------------------------------------------------------------
# Text extraction — multiple file types
# ---------------------------------------------------------------------------

def extract_text_from_file(file_bytes, filename):
    """Returns list of (page_number, text) tuples."""
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    pages = []

    if ext == 'pdf':
        reader = PdfReader(BytesIO(file_bytes))
        for i, page in enumerate(reader.pages, 1):
            pages.append((i, page.extract_text() or ""))

    elif ext == 'txt':
        text = file_bytes.decode('utf-8', errors='replace')
        lines = text.splitlines()
        for i in range(0, len(lines), 50):
            pages.append((i // 50 + 1, "\n".join(lines[i:i + 50])))

    elif ext == 'csv':
        text = file_bytes.decode('utf-8', errors='replace')
        rows = list(csv.reader(StringIO(text)))
        for i in range(0, len(rows), 50):
            block = "\n".join([",".join(r) for r in rows[i:i + 50]])
            pages.append((i // 50 + 1, block))

    elif ext == 'docx' and HAS_DOCX:
        doc = DocxDocument(BytesIO(file_bytes))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        for i in range(0, len(paras), 30):
            pages.append((i // 30 + 1, "\n".join(paras[i:i + 30])))

    elif ext == 'pptx' and HAS_PPTX:
        prs = Presentation(BytesIO(file_bytes))
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            pages.append((i, "\n".join(texts)))

    elif ext in ('xlsx', 'xls') and HAS_XLSX:
        wb = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True)
        page_num = 1
        for sheet in wb.worksheets:
            rows = ["\t".join([str(c) if c is not None else "" for c in row])
                    for row in sheet.iter_rows(values_only=True)]
            for i in range(0, len(rows), 50):
                pages.append((page_num, "\n".join(rows[i:i + 50])))
                page_num += 1
    else:
        # Generic fallback so "any filetype" actually works: anything that
        # isn't one of the formats above (md, json, xml, html, log, py, js,
        # yaml, code files, etc.) is treated as plain text if it decodes.
        # Only genuinely binary, non-text files (images, zips, exe...) fail.
        try:
            text = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            try:
                text = file_bytes.decode('latin-1')
            except Exception:
                raise ValueError(
                    f"Could not read .{ext or '(no extension)'} — it looks like a "
                    f"binary format that isn't text-extractable (e.g. an image, "
                    f"archive, or executable)."
                )
        lines = text.splitlines()
        for i in range(0, len(lines), 50):
            pages.append((i // 50 + 1, "\n".join(lines[i:i + 50])))

    return pages


def clean_text(text):
    """Normalize text — remove OCR artifacts and extra whitespace."""
    text = " ".join(text.split())
    text = re.sub(r'(\b\w+\b)(\s+\1){2,}', r'\1', text)  # remove repeated words
    return text


def chunk_pages(pages):
    chunks = []
    idx = 1
    for page_no, text in pages:
        text = clean_text(text)
        if not text:
            continue
        for i in range(0, len(text), CHUNK_SIZE):
            chunks.append({"label": f"c{idx}", "page": page_no, "text": text[i:i + CHUNK_SIZE]})
            idx += 1
    return chunks


# ---------------------------------------------------------------------------
# Chapter detection — the model reads a short preview of every page and
# proposes where each chapter starts, so a 400-page book doesn't have to be
# split by hand. Falls back to a single "Full Document" chapter if detection
# fails or the book has no headings the model can find in the extracted text
# (e.g. a chapter title baked into an image rather than real text).
# ---------------------------------------------------------------------------

CHAPTER_PREVIEW_CHARS = 400


def detect_chapters(pages, on_status=None):
    if not pages:
        return [{"number": 1, "title": "Full Document", "start_page": 1, "end_page": 1}]

    last_page = pages[-1][0]
    preview_blob = "\n".join(
        f"--- page {page_no} ---\n{(text or '').strip()[:CHAPTER_PREVIEW_CHARS]}"
        for page_no, text in pages
    )

    prompt = (
        "Below are the first characters of every page of a textbook, in order. "
        "Identify where each chapter begins. Chapter headings usually appear near the "
        "top of a page — as a chapter number, the word 'Chapter', or a distinct title "
        "line. Use your judgement based on the text pattern across pages.\n\n"
        "Respond with ONLY valid JSON, no markdown fences and no commentary: a list of "
        "objects with keys \"number\" (integer, sequential starting at 1), \"title\" "
        "(string), and \"start_page\" (integer — the page this chapter begins on). "
        "If there is meaningful front matter (cover, preface, table of contents) before "
        "chapter 1, include it as number 0 titled \"Front Matter\". List chapters in "
        "page order.\n\n"
        f"{preview_blob}"
    )

    try:
        model = genai.GenerativeModel(MODEL)
        response = None
        for attempt in range(4):
            try:
                response = model.generate_content(prompt)
                break
            except Exception as e:
                if ('429' in str(e) or 'rate' in str(e).lower()) and attempt < 3:
                    wait = 20 + attempt * 15
                    if on_status:
                        on_status(f"Rate limited while detecting chapters — waiting {wait}s...")
                    time.sleep(wait)
                else:
                    raise

        raw = (response.text or "").strip()
        raw = re.sub(r"^```(json)?|```$", "", raw, flags=re.MULTILINE).strip()
        data = json.loads(raw)
        if not isinstance(data, list) or not data:
            raise ValueError("model returned no chapters")

        chapters = []
        for item in data:
            num = int(item["number"])
            title = str(item.get("title") or f"Chapter {num}").strip()
            start = max(1, min(int(item["start_page"]), last_page))
            chapters.append({"number": num, "title": title, "start_page": start})

        chapters.sort(key=lambda c: c["start_page"])
        deduped, seen_starts = [], set()
        for c in chapters:
            if c["start_page"] in seen_starts:
                continue
            seen_starts.add(c["start_page"])
            deduped.append(c)
        chapters = deduped

        if not chapters or chapters[0]["start_page"] != 1:
            chapters.insert(0, {
                "number": 0, "title": "Front Matter", "start_page": 1,
            })
            chapters.sort(key=lambda c: c["start_page"])

        for i, c in enumerate(chapters):
            c["end_page"] = (chapters[i + 1]["start_page"] - 1) if i + 1 < len(chapters) else last_page

        return chapters

    except Exception as e:
        if on_status:
            on_status(f"⚠️ Automatic chapter detection failed ({e}) — treating document as one chapter.")
        return [{"number": 1, "title": "Full Document", "start_page": 1, "end_page": last_page}]


def _chapter_for_page(chapters, page_no):
    for c in chapters:
        if c["start_page"] <= page_no <= c["end_page"]:
            return c
    return chapters[-1]


def chunk_pages_by_chapter(pages, chapters):
    chunks = []
    idx = 1
    for page_no, text in pages:
        text = clean_text(text)
        if not text:
            continue
        chapter = _chapter_for_page(chapters, page_no)
        for i in range(0, len(text), CHUNK_SIZE):
            chunks.append({
                "label": f"c{idx}",
                "page": page_no,
                "text": text[i:i + CHUNK_SIZE],
                "chapter_number": chapter["number"],
                "chapter_title": chapter["title"],
            })
            idx += 1
    return chunks


# ---------------------------------------------------------------------------
# Embeddings (Google free API)
# ---------------------------------------------------------------------------

def generate_embeddings(texts, on_progress=None, on_status=None, max_retries=6):
    """
    Embeds texts in batches with exponential backoff on rate limits.
    on_progress(done_count, total_count) is called after each successful batch.
    on_status(message) is called during retries so the caller can surface
    "rate limited, retrying..." to the UI instead of it looking frozen.
    Gives up after max_retries on a single batch and raises, rather than
    retrying forever (which is what happens if a *daily* quota is exhausted —
    every retry just hits 429 again, indefinitely).
    """
    all_embeddings = []
    total = len(texts)
    for i in range(0, total, EMBED_BATCH):
        batch = texts[i:i + EMBED_BATCH]
        attempt = 0
        while True:
            try:
                result = genai.embed_content(
                    model=EMBED_MODEL,
                    content=batch,
                    task_type="retrieval_document",
                )
                all_embeddings.extend(result['embedding'])
                break
            except Exception as e:
                if ('429' in str(e) or 'rate' in str(e).lower()) and attempt < max_retries:
                    attempt += 1
                    wait = min(60, (2 ** attempt) + random.uniform(0, 1))
                    msg = f"Rate limited, retrying in {wait:.0f}s (attempt {attempt}/{max_retries})..."
                    print(f"[batch {i}] {msg}")
                    if on_status:
                        on_status(msg)
                    time.sleep(wait)
                elif '429' in str(e) or 'rate' in str(e).lower():
                    raise RuntimeError(
                        f"Gemini embedding API kept rate-limiting after {max_retries} retries. "
                        f"This usually means your daily free-tier quota is exhausted, not just a "
                        f"per-minute limit — check https://aistudio.google.com/app/apikey for usage, "
                        f"or wait and try again later. Last error: {e}"
                    )
                else:
                    raise
        if on_progress:
            on_progress(min(i + EMBED_BATCH, total), total)
    return all_embeddings


def embed_query(text):
    result = genai.embed_content(
        model=EMBED_MODEL,
        content=text,
        task_type="retrieval_query",
    )
    return result['embedding']


# ---------------------------------------------------------------------------
# Routes: documents
# ---------------------------------------------------------------------------

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/documents", methods=["GET"])
def list_documents():
    conn = get_db()
    try:
        cur = conn.cursor()
        cur.execute(
            "SELECT d.DocumentId, d.FileName, d.PageCount, d.ChunkCount, d.UploadedAt, "
            "(SELECT COUNT(DISTINCT ChapterNumber) FROM dbo.Chunks "
            " WHERE DocumentId = d.DocumentId AND ChapterNumber IS NOT NULL AND ChapterNumber != 0) AS ChapterCount "
            "FROM dbo.Documents d ORDER BY d.UploadedAt DESC"
        )
        rows = cur.fetchall()
        return jsonify([
            {"id": r.DocumentId, "name": r.FileName, "pages": r.PageCount,
             "chunks": r.ChunkCount, "chapters": r.ChapterCount or 0,
             "uploaded_at": r.UploadedAt.isoformat()}
            for r in rows
        ])
    finally:
        conn.close()


@app.route("/api/documents/<int:doc_id>", methods=["DELETE"])
def delete_document(doc_id):
    conn = get_db()
    try:
        cur = conn.cursor()
        cur.execute("DELETE FROM dbo.Documents WHERE DocumentId = ?", doc_id)
        conn.commit()
        return jsonify({"ok": True})
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# Upload jobs — background processing with progress polling.
#
# Large documents (hundreds of pages -> thousands of chunks -> tens of
# embedding API calls, some rate-limited) can take many minutes. Doing that
# inside a single blocking HTTP request is what was breaking 400-page
# uploads: the request looked "stuck" with no real feedback and could be
# killed by any timeout in the chain. Instead we kick off a background
# thread immediately, return a job_id, and let the frontend poll for
# real progress (extracting -> chunking -> embedding X/Y -> saving -> done).
# ---------------------------------------------------------------------------

JOBS = {}
JOBS_LOCK = threading.Lock()


def _set_job(job_id, **fields):
    with JOBS_LOCK:
        JOBS[job_id].update(fields)


def _process_upload_job(job_id, file_bytes, filename):
    try:
        _set_job(job_id, stage="extracting", message=f'Extracting text from "{filename}"...')
        pages = extract_text_from_file(file_bytes, filename)

        chapters = [{"number": 1, "title": "Full Document", "start_page": 1,
                     "end_page": pages[-1][0] if pages else 1}]
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        if ext == 'pdf' and len(pages) > 3:
            # Only worth an extra model call for documents long enough to
            # plausibly have chapters — a 1-3 page file is not a textbook.
            _set_job(job_id, stage="detecting_chapters", message="Detecting chapter structure...")
            chapters = detect_chapters(pages, on_status=lambda msg: _set_job(job_id, message=msg))

        chapter_count = len([c for c in chapters if c["number"] != 0])
        _set_job(job_id, stage="chunking",
                 message=f"Found {chapter_count} chapter(s) — splitting into chunks..."
                         if chapter_count > 1 else "Splitting into chunks...")
        chunks = chunk_pages_by_chapter(pages, chapters)
        if not chunks:
            _set_job(job_id, stage="error", error="No extractable text found in this file.")
            return

        total_chunks = len(chunks)
        _set_job(
            job_id, stage="embedding", done=0, total=total_chunks,
            message=f"Generating embeddings 0/{total_chunks}...",
        )

        def on_progress(done, total):
            _set_job(job_id, done=done, total=total,
                        message=f"Generating embeddings {done}/{total}...")                                

        def on_status(msg):
            _set_job(job_id, message=msg)   

        texts = [c["text"] for c in chunks]
        embed_skipped = False                       
        try:
            embeddings = generate_embeddings(texts, on_progress=on_progress, on_status=on_status)
            for c, emb in zip(chunks, embeddings):
                c["embedding"] = json.dumps(emb)
        except Exception as emb_err:
            # Quota exhausted or any other embedding failure:
            # mark chunks with no embedding and continue — the document
            # still uploads fine, keyword search still works, and
            # semantic_search returns a clear "no embeddings" message.
            for c in chunks:
                c["embedding"] = None
            embed_skipped = True
            _set_job(job_id, message=f"⚠️ Embeddings skipped ({emb_err}). Keyword search still works.")

        _set_job(job_id, stage="saving", message=f"Saving {total_chunks} chunks to SQL Server...")
        conn = get_db()
        try:
            cur = conn.cursor()
            cur.fast_executemany = True
            cur.execute(
                "INSERT INTO dbo.Documents (FileName, PageCount, ChunkCount) "
                "OUTPUT INSERTED.DocumentId VALUES (?, ?, ?)",
                filename, len(pages), len(chunks),
            )
            doc_id = cur.fetchone()[0]
            cur.executemany(
                "INSERT INTO dbo.Chunks (DocumentId, ChunkLabel, PageNumber, ChunkText, Embedding, ChapterNumber, ChapterTitle) "
                "VALUES (?, ?, ?, ?, ?, ?, ?)",
                [(doc_id, c["label"], c["page"], c["text"], c.get("embedding"),
                  c["chapter_number"], c["chapter_title"]) for c in chunks],
            )
            conn.commit()
        except Exception as e:
            conn.rollback()
            _set_job(job_id, stage="error", error=str(e))
            return
        finally:
            conn.close()

        _set_job(
            job_id, stage="done", message="Done." if not embed_skipped else "Done (no embeddings — semantic search disabled).",
            result={"id": doc_id, "name": filename, "pages": len(pages), "chunks": len(chunks),
                    "chapters": chapter_count, "embed_skipped": embed_skipped},
        )
    except Exception as e:
        _set_job(job_id, stage="error", error=str(e))


@app.route("/api/upload", methods=["POST"])
def upload():
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "No file uploaded"}), 400

    filename = file.filename
    file_bytes = file.read()
    # No extension whitelist: extract_text_from_file has a generic
    # text-decode fallback, so any text-based filetype is accepted.
    # It will still fail gracefully for genuinely binary files.

    job_id = str(uuid.uuid4())
    with JOBS_LOCK:
        JOBS[job_id] = {"stage": "queued", "message": "Queued..."}

    thread = threading.Thread(
        target=_process_upload_job, args=(job_id, file_bytes, filename), daemon=True
    )
    thread.start()

    return jsonify({"job_id": job_id})

@app.route("/api/upload/status/<job_id>", methods=["GET"])
def upload_status(job_id):
    with JOBS_LOCK:
        job = JOBS.get(job_id)
    if not job:
        return jsonify({"error": "Unknown job_id"}), 404
    return jsonify(job)


# ---------------------------------------------------------------------------
# Tool implementations
# ---------------------------------------------------------------------------

def tool_get_index(conn, doc_id):
    cur = conn.cursor()
    cur.execute(
        "SELECT ChapterNumber, ChapterTitle, MIN(PageNumber) AS StartPage, "
        "MAX(PageNumber) AS EndPage, COUNT(*) AS ChunkCount "
        "FROM dbo.Chunks WHERE DocumentId = ? GROUP BY ChapterNumber, ChapterTitle",
        doc_id,
    )
    rows = cur.fetchall()
    has_chapters = any(r.ChapterNumber is not None for r in rows)

    if not has_chapters:
        # Legacy document uploaded before chapter detection existed.
        cur.execute(
            "SELECT ChunkLabel, PageNumber, LEFT(ChunkText, 90) AS Preview "
            "FROM dbo.Chunks WHERE DocumentId = ? ORDER BY ChunkId", doc_id,
        )
        return [{"id": r.ChunkLabel, "page": r.PageNumber, "preview": r.Preview}
                for r in cur.fetchall()]

    chapters = sorted(rows, key=lambda r: (r.ChapterNumber if r.ChapterNumber is not None else 0))
    return {
        "chapters": [
            {"number": r.ChapterNumber, "title": r.ChapterTitle,
             "pages": f"{r.StartPage}-{r.EndPage}", "chunk_count": r.ChunkCount}
            for r in chapters
        ],
        "note": "Call get_chapter_chunks(chapter_number) to see the chunk ids inside a "
                "chapter, then get_chunk to read one.",
    }


def tool_get_chapter_chunks(conn, doc_id, chapter_number):
    cur = conn.cursor()
    cur.execute(
        "SELECT ChunkLabel, PageNumber, LEFT(ChunkText, 90) AS Preview "
        "FROM dbo.Chunks WHERE DocumentId = ? AND ChapterNumber = ? ORDER BY ChunkId",
        doc_id, chapter_number,
    )
    rows = cur.fetchall()
    if not rows:
        return {"error": f"No chunks found for chapter {chapter_number}. Check get_index for valid chapter numbers."}
    return [{"id": r.ChunkLabel, "page": r.PageNumber, "preview": r.Preview} for r in rows]


def tool_get_chunk(conn, doc_id, chunk_label):
    cur = conn.cursor()
    cur.execute(
        "SELECT ChunkLabel, PageNumber, ChunkText FROM dbo.Chunks "
        "WHERE DocumentId = ? AND ChunkLabel = ?", doc_id, chunk_label,
    )
    r = cur.fetchone()
    if not r:
        return {"error": f"No chunk with id {chunk_label}"}
    return {"id": r.ChunkLabel, "page": r.PageNumber, "text": r.ChunkText}


def tool_search_chunks(conn, doc_id, keyword):
    """Improved keyword search — tries multiple normalizations until results found."""
    cur = conn.cursor()
    variants = [
        keyword,
        keyword.upper(),
        keyword.lower(),
        re.sub(r"[^a-zA-Z0-9 ]", "", keyword),   # strip punctuation
        keyword.split()[0] if keyword.split() else keyword,  # first word only
    ]
    seen_ids = set()
    results = []
    for v in variants:
        if not v.strip():
            continue
        cur.execute(
            "SELECT TOP 15 ChunkLabel, PageNumber, LEFT(ChunkText, 90) AS Preview "
            "FROM dbo.Chunks WHERE DocumentId = ? AND LOWER(ChunkText) LIKE LOWER(?) "
            "ORDER BY ChunkId",
            doc_id, f"%{v}%",
        )
        for r in cur.fetchall():
            if r.ChunkLabel not in seen_ids:
                seen_ids.add(r.ChunkLabel)
                results.append({"id": r.ChunkLabel, "page": r.PageNumber, "preview": r.Preview})
        if results:
            break  # stop at first successful variant
    return results[:15]


def tool_semantic_search(conn, doc_id, query, top_k=5):
    """FAISS vector similarity search using stored embeddings."""
    try:
        cur = conn.cursor()
        cur.execute(
            "SELECT ChunkLabel, PageNumber, LEFT(ChunkText, 90) AS Preview, Embedding "
            "FROM dbo.Chunks WHERE DocumentId = ? AND Embedding IS NOT NULL ORDER BY ChunkId",
            doc_id,
        )
        rows = cur.fetchall()
        if not rows:
            return {"error": "No embeddings stored. Re-upload the document to generate embeddings."}

        labels = [r.ChunkLabel for r in rows]
        pages = [r.PageNumber for r in rows]
        previews = [r.Preview for r in rows]
        vectors = np.array([json.loads(r.Embedding) for r in rows], dtype=np.float32)

        # Build FAISS flat index and search
        dim = vectors.shape[1]
        index = faiss.IndexFlatL2(dim)
        index.add(vectors)

        q_vec = np.array([embed_query(query)], dtype=np.float32)
        k = min(top_k, len(labels))
        distances, indices = index.search(q_vec, k)

        results = []
        for dist, idx in zip(distances[0], indices[0]):
            if idx < 0:
                continue
            results.append({
                "id": labels[idx],
                "page": pages[idx],
                "preview": previews[idx],
                "score": round(float(dist), 4),
            })
        return results

    except Exception as e:
        return {"error": f"Semantic search failed: {e}"}


def run_tool(conn, doc_id, name, tool_input):
    if name == "get_index":
        return tool_get_index(conn, doc_id)
    if name == "get_chapter_chunks":
        return tool_get_chapter_chunks(conn, doc_id, tool_input.get("chapter_number"))
    if name == "get_chunk":
        return tool_get_chunk(conn, doc_id, tool_input.get("chunk_id", ""))
    if name == "search_chunks":
        return tool_search_chunks(conn, doc_id, tool_input.get("keyword", ""))
    if name == "semantic_search":
        return tool_semantic_search(conn, doc_id, tool_input.get("query", ""), tool_input.get("top_k", 5))
    return {"error": f"unknown tool {name}"}


# ---------------------------------------------------------------------------
# /api/ask — streams every AI<->Code exchange via SSE
# ---------------------------------------------------------------------------

@app.route("/api/ask", methods=["POST"])
def ask():
    payload = request.get_json(force=True)
    doc_id = payload.get("document_id")
    question = (payload.get("question") or "").strip()

    if not doc_id or not question:
        return jsonify({"error": "document_id and question are required"}), 400

    def event(event_type, **data):
        return f"data: {json.dumps({'type': event_type, **data})}\n\n"

    def generate():
        conn = get_db()
        try:
            yield event("question", text=question)
            model = genai.GenerativeModel(
                model_name=MODEL,
                system_instruction=SYSTEM_PROMPT,
                tools=TOOLS,
                generation_config=genai.types.GenerationConfig(
                    max_output_tokens=8192,
                ),
            )
            chat = model.start_chat(enable_automatic_function_calling=False)
            pending_message = question

            for turn in range(MAX_TURNS):
                yield event("thinking_start")
                response = None
                for attempt in range(5):
                    try:
                        response = chat.send_message(pending_message)
                        break
                    except Exception as e:
                        if ('429' in str(e) or 'rate' in str(e).lower()) and attempt < 4:
                            wait = 30 + attempt * 15  # 30s, 45s, 60s, 75s
                            yield event("ai_thought", text=f"⏳ Rate limited — waiting {wait}s before retrying (attempt {attempt+1}/5)...")
                            time.sleep(wait)
                        else:
                            yield event("dead_end", text=f"Request to the model failed: {e}")
                            return
                if response is None:
                    yield event("dead_end", text="Model did not respond after retries.")
                    return

                try:
                    parts = response.candidates[0].content.parts
                except Exception:
                    yield event("dead_end", text="Model returned an empty response.")
                    return

                text_parts = "".join(p.text for p in parts if getattr(p, "text", None)).strip()
                function_calls = [p.function_call for p in parts if getattr(p, "function_call", None)]

                if text_parts and not function_calls:
                    yield event("final_answer", text=text_parts)
                    return  # kill switch

                if text_parts and function_calls:
                    yield event("ai_thought", text=text_parts)

                if not function_calls:
                    yield event("dead_end", text="Model stopped without an answer or a tool call.")
                    return

                response_parts = []
                for call in function_calls:
                    tool_input = dict(call.args) if call.args else {}
                    yield event("ai_call", tool=call.name, input=tool_input)
                    result = run_tool(conn, doc_id, call.name, tool_input)

                    if isinstance(result, list) and len(result) == 0:
                        # Don't dead-end — tell the model nothing was found
                        # so it can try a different keyword or approach
                        result = {"message": "No chunks matched this search. Try a different keyword or broader term."}
                    else:
                        yield event("code_response", tool=call.name, result=result)

                    response_parts.append(
                        genai.protos.Part(
                            function_response=genai.protos.FunctionResponse(
                                name=call.name,
                                response={"result": json.dumps(result)},
                            )
                        )
                    )

                pending_message = genai.protos.Content(parts=response_parts)

            yield event("dead_end", text="Stopped after maximum reasoning steps without a confident answer.")
        finally:
            conn.close()

    return Response(stream_with_context(generate()), mimetype="text/event-stream")


if __name__ == "__main__":
    # threaded=True is required: without it Flask's dev server handles one
    # request at a time, so the /api/upload/status polling requests would
    # queue up behind the background upload job instead of running alongside it.
    app.run(debug=True, port=5000, threaded=True)